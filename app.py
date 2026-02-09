import streamlit as st
from pptx import Presentation
from pptx.util import Pt
import pdfplumber
import io

# Configuración de la aplicación
st.set_page_config(page_title="MT Valero - Editor Pro", layout="wide")

# Memoria de sesión
if 'hojas' not in st.session_state:
    st.session_state.hojas = []
if 'datos_pdf' not in st.session_state:
    st.session_state.datos_pdf = {"cliente": "", "fecha": ""}

st.title("📋 Generador de Reportes: Inyección en Tablas y Cuadros")

# --- BARRA LATERAL ---
with st.sidebar:
    st.header("1. Carga de Recursos")
    plantilla = st.file_uploader("Subir tu Plantilla (.pptx)", type=["pptx"])
    archivo_pdf = st.file_uploader("Subir Hoja Valero (PDF)", type=["pdf"])
    fotos_totales = st.file_uploader("Galería de WhatsApp", type=["jpg", "png", "jpeg"], accept_multiple_files=True)
    
    if archivo_pdf and st.button("🔍 Escanear PDF de Valero"):
        with pdfplumber.open(archivo_pdf) as pdf:
            texto = pdf.pages[0].extract_text()
            for l in texto.split('\n'):
                if "Cliente:" in l: st.session_state.datos_pdf["cliente"] = l.split("Cliente:")[1].strip()
                if "Fecha:" in l: st.session_state.datos_pdf["fecha"] = l.split("Fecha:")[1].strip()
        st.success("Datos extraídos correctamente.")

# --- CUERPO PRINCIPAL ---
if not plantilla:
    st.info("👈 Por favor, sube tu archivo 'Reporte_MT_Final.pptx' para comenzar.")
else:
    prs_base = Presentation(io.BytesIO(plantilla.getvalue()))
    nombres_layouts = [layout.name for layout in prs_base.slide_layouts]

    col_edit, col_prev = st.columns([1, 1])

    with col_edit:
        st.subheader("📝 Configurar Diapositiva")
        with st.form("editor_hoja", clear_on_submit=True):
            diseno = st.selectbox("Elegir Diseño (Nombre en PPT):", nombres_layouts)
            texto_tecnico = st.text_area("Texto para la Tabla o Cuadro (Times New Roman 12):")
            
            st.write("🖼️ **Fotos para los cuadros predeterminados:**")
            fotos_escogidas = []
            if fotos_totales:
                cols_img = st.columns(3)
                for i, foto in enumerate(fotos_totales):
                    with cols_img[i % 3]:
                        st.image(foto, width=90)
                        if st.checkbox("Usar", key=f"sel_{foto.name}_{len(st.session_state.hojas)}"):
                            fotos_escogidas.append(foto)
            
            if st.form_submit_button("➕ GUARDAR DIAPOSITIVA"):
                st.session_state.hojas.append({
                    "layout_idx": nombres_layouts.index(diseno),
                    "nombre": diseno,
                    "texto": texto_tecnico,
                    "fotos": fotos_escogidas
                })
                st.rerun()

    with col_prev:
        st.subheader("👁️ Vista Previa")
        if st.session_state.hojas:
            num = st.number_input("Ver Hoja #:", min_value=1, max_value=len(st.session_state.hojas)) - 1
            h = st.session_state.hojas[num]
            with st.container(border=True):
                st.markdown(f"**Diseño:** {h['nombre']}")
                st.markdown(f"**Texto Guardado:** {h['texto'][:100]}...")
                st.write(f"🖼️ Fotos: {len(h['fotos'])}")
            
            if st.button("🗑️ Eliminar Hoja"):
                st.session_state.hojas.pop(num)
                st.rerun()

    # --- GENERACIÓN DEL ARCHIVO FINAL ---
    st.divider()
    if st.session_state.hojas:
        if st.button("🚀 GENERAR Y DESCARGAR REPORTE FINAL"):
            prs_final = Presentation(io.BytesIO(plantilla.getvalue()))
            
            # Borramos las diapositivas de ejemplo que vienen en el archivo original
            for i in range(len(prs_final.slides) - 1, -1, -1):
                rId = prs_final.slides._sldIdLst[i].rId
                prs_final.part.drop_rel(rId)
                del prs_final.slides._sldIdLst[i]

            for h in st.session_state.hojas:
                slide = prs_final.slides.add_slide(prs_final.slide_layouts[h['layout_idx']])
                
                # 1. BUSCAR LA TABLA O CUADRO DE TEXTO
                img_ptr = 0
                for shape in slide.shapes:
                    # Caso A: Si es una Tabla
                    if shape.has_table:
                        # Buscamos la celda de descripción (usualmente la última o vacía)
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text == "" or "DESCRIPCIÓN" in cell.text.upper():
                                    cell.text = h['texto']
                                    # Aplicamos formato Times New Roman 12
                                    for paragraph in cell.text_frame.paragraphs:
                                        for run in paragraph.runs:
                                            run.font.name = 'Times New Roman'
                                            run.font.size = Pt(12)
                    
                    # Caso B: Si es un Placeholder de Texto (Cuadro de texto predeterminado)
                    elif shape.is_placeholder and shape.placeholder_format.type in [2, 7]:
                        shape.text = h['texto']
                        for p in shape.text_frame.paragraphs:
                            for r in p.runs:
                                r.font.name = 'Times New Roman'
                                r.font.size = Pt(12)

                    # Caso C: Si es un Placeholder de Imagen (Cuadro de foto)
                    elif shape.is_placeholder and (shape.placeholder_format.type == 18 or "PICTURE" in shape.name.upper()):
                        if img_ptr < len(h['fotos']):
                            shape.insert_picture(io.BytesIO(h['fotos'][img_ptr].read()))
                            h['fotos'][img_ptr].seek(0)
                            img_ptr += 1

            output = io.BytesIO()
            prs_final.save(output)
            output.seek(0)
            
            st.download_button(
                label="📥 DESCARGAR REPORTE PPTX EDITADO",
                data=output.getvalue(),
                file_name="Reporte_Final_MT_Editado.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )